import json

import pandas as pd
import requests
from uuid import uuid4
from django.http import HttpResponse
from openpyxl import Workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from rest_framework import status
from rest_framework.response import Response
from rest_framework.views import APIView

from datetime import date

class FetchAllChartsData(APIView):

    """Fetching all data from charts api's"""

    def get(self, request, account_name):
        # Geting all data from api's using account name
        try:

            list_urls = {
                "performance": "v1",
                "comparison": "v1",
                "transactions": "v2",
                "deals": "v2",
                "latestnews": "v1",
                "offers": "v2",
            }
            wb = Workbook()
            # grab the active worksheet
            ws = wb.active
            ws.append(["slide_no", "object_type", "content", "font_size", "color"])
            
            title_name = "Account Insights Books –" + account_name
            ws.append([0, "Title 4", title_name, 44, "0,0,0"])
            today = date.today()
            ws.append([0, "TextBox 1", date.isoformat(today), 14, "0,0,0"])

            for slide in range(1, 13):
                title_name = "Account strategy and planning –" + account_name
                ws.append([slide, "Title 1", title_name, 26, "0,0,0"])
                # ws.append([slide,'Subtitle 3','Test',18,'0xFB,0x8F,0x00'])

            for one_url, value in list_urls.items():
                if value == "v1":
                    url = "http://0.0.0.0:8000/api/" + account_name + "/" + one_url
                elif value == "v2":
                    url = "http://0.0.0.0:8000/api/v2/" + account_name + "/" + one_url
                else:
                    pass
                print(url, "url")
                try:
                    response = requests.get(url)
                    if response.status_code == 200:
                        data = json.loads(json.dumps(response.json()))
                    else:
                        continue
                except:
                    return Response(
                        data={"msg": "Not able fetch data from external api's"},
                        content_type="application/json",
                        status=status.HTTP_200_OK,
                    )
                if one_url == "performance":
                    # REvenue data set
                    try:
                        final_data = {}
                        category = []
                        series = {}
                        for key, value in data["revenue_by_segment"].items():
                            sublist = []
                            for keyvalue in value.keys():
                                category.append(value.keys())
                                if "value" in value[keyvalue]["segment_revenues_all"].keys():
                                    sublist.append(value[keyvalue]["segment_revenues_all"]["value"])
                                else:
                                    sublist.append(0)
                            series[key] = sublist
                        final_data["categories"] = list(category[0])
                        final_data["series"] = series
                        ws.append([1, "Chart 13", json.dumps(final_data), "", ""])
                    except Exception as e:
                        print("Exception while revenue data", e)
                    # MArket captes graph data set
                    try:
                        final_data = {}
                        final_data["categories"] = list(data["marketcap"].keys())
                        final_data["series"] = {"series1": list(data["marketcap"].values())}
                        ws.append([1, "Chart 14", json.dumps(final_data), "", ""])
                    except Exception as e:
                        print("Exception while market captes data ", e)
                    # ebitda dataset
                    try:
                        final_data = {}
                        final_data["categories"] = list(data["ebitda"].keys())
                        series = []
                        for onerecord in data["ebitda"].keys():
                            if len(data["ebitda"][onerecord]["ebitda_adj_ind"].keys()) > 0:
                                series.append(data["ebitda"][onerecord]["ebitda_adj_ind"]["value"])
                            else:
                                series.append(0)
                        final_data["series"] = {"series1": series}
                        ws.append([1, "Chart 15", json.dumps(final_data), "", ""])
                    except Exception as e:
                        print("Exception while ebitda   data", e)
                elif one_url == "comparison":
                    try:
                        for key in data.keys():
                            for onerow in data[key]:
                                final_data = {}
                                final_data["categories"] = list(data[key][onerow].keys())
                                final_data["series"] = {"series1":[value if key in ['fte'] else value*100 for value in list(data[key][onerow].values())]}
                                if onerow == "yoy":
                                    ws.append([2, "Chart 21", json.dumps(final_data), "", ""])
                                elif onerow == "3_years":
                                    ws.append([2, "Chart 22", json.dumps(final_data), "", ""])
                                elif onerow == "5_years":
                                    ws.append([2, "Chart 23", json.dumps(final_data), "", ""])
                                elif onerow == "adjusted_revenue":
                                    ws.append([3, "Chart 21", json.dumps(final_data), "", ""])
                                elif onerow == "total_assets":
                                    ws.append([3, "Chart 22", json.dumps(final_data), "", ""])
                                elif onerow == "fte":
                                    ws.append([3, "Chart 23", json.dumps(final_data), "", ""])
                                elif onerow == "asset_growth":
                                    ws.append([4, "Chart 21", json.dumps(final_data), "", ""])
                                elif onerow == "fte_growth":
                                    ws.append([4, "Chart 22", json.dumps(final_data), "", ""])
                                elif onerow == "revenue":
                                    ws.append([4, "Chart 23", json.dumps(final_data), "", ""])
                                else:
                                    pass
                    except Exception as e:
                        print("Exception while comparison data", e)
                elif one_url == "transactions":
                    try:

                        headers = [
                            "Buyer company",
                            "Target company",
                            "Announce date",
                            "Closing date",
                            # "Comments",
                            "Market Cap (in USDm)",
                            "Intensity",
                        ]
                        final_data = []
                        final_data.append(headers)
                        for record in data:
                            innerlist = []
                            for req_value in [
                                "buyer_company_name",
                                "target_company_name",
                                "announced_date",
                                "closing_date",
                                # "comments",
                                "market_cap",
                                "intensity",
                            ]:
                                if req_value == "comments":
                                    if len(record[req_value]) > 200:
                                        innerlist.append(str(record[req_value][:200]))
                                    else:
                                        innerlist.append(str(record[req_value]))
                                else:
                                    innerlist.append(str(record[req_value]))
                            # print(innerlist,"innerlistinnerlistinnerlist")
                            final_data.append(innerlist)
                        ws.append([6, "Table 5", str(final_data), "", ""])
                    except Exception as e:
                        print("Exception while transactions data", e)
                elif one_url == "deals":
                    try:
                        if "past" in data.keys():
                            headers = ["Vendor", "Year", "Theme", "TCV ($M)", "Contract End Period"]
                            final_data = []
                            final_data.append(headers)
                            for record in data["past"]:
                                innerlist = []
                                for requiredkey in ["vendor_name", "year", "theme", "tcv", "end_renewal"]:
                                    innerlist.append(str(record[requiredkey]))
                                final_data.append(innerlist)
                            ws.append([11, "Table 5", json.dumps(final_data), "", ""])
                        if "deal_by_peer" in data.keys():
                            headers = ["Theme", "Vendor", "ITS Provider", "TCV ($M)", "Year of deal", "Duration"]
                            final_data = []
                            final_data.append(headers)
                            for record in data["deal_by_peer"]:
                                innerlist = []
                                for requiredkey in [
                                    "theme",
                                    "client_name",
                                    "vendor_name",
                                    "tcv",
                                    "announcement_date",
                                    "duration",
                                ]:
                                    innerlist.append(str(record[requiredkey]))
                                final_data.append(innerlist)
                            ws.append([12, "Table 5", json.dumps(final_data), "", ""])
                    except Exception as e:
                        print("Exception while deals data", e)

                elif one_url == "latestnews":
                    try:
                        # Latest news data append to excel
                        headers = ["Developments/news", "Date", "Link"]
                        final_data = []
                        final_data.append(headers)
                        for onerecord in data:
                            final_data.append(list(onerecord.values()))
                        ws.append([5, "Table 5", json.dumps(final_data), "", ""])
                    except Exception as e:
                        print("Exception while latestnews data", e)

                elif one_url == "offers":
                    try:
                        final_data = []
                        usecase = []
                        impact_metrics = []
                        skills_completencie = []
                        exports = {}
                        mcK_experts={}
                        skills_completencies_dict = {}
                        for onerecord in data["industry_specific"]:
                            usecase.append(onerecord["use_cases"])
                            impact_metrics.append(onerecord["impact_metrics"])
                            if onerecord["skills_completencies"] in onerecord.keys():
                                skills_completencie.extend(
                                    [skill for skill in onerecord["skills_completencies"].split("|") if skill not in skills_completencies]
                                )
                            exports = {
                                export.split(",", 1)[0]: export.split(",", 1)[1:]
                                for export in onerecord["experts"].split("|")
                            }

                            mcK_experts = {
                                export.split(",", 1)[0]: export.split(",", 1)[1:]
                                for export in onerecord["mcK_experts"].split("|")
                            }

                            oneDict = {}
                            oneDict["offer"] = onerecord["offer"]
                            oneDict["value"] = onerecord["value"]
                            oneDict["code"] = onerecord["code"]
                            oneDict["percent"] = onerecord["percent"]
                            final_data.append(oneDict)
                        skills_completencies_dict["Skills and Competencies"] = skills_completencie
                        ws.append([8, "Rectangle 6", json.dumps(final_data), "", ""])
                        ws.append([9, "Text 1", json.dumps(list(set(usecase))), "", ""])
                        ws.append([9, "Text 4", json.dumps(exports), "", ""])
                        ws.append([9, "Text 5", json.dumps(mcK_experts), "", ""])

                        ws.append([10, "Text 1", json.dumps(list(set(impact_metrics))), "", ""])
                        ws.append([10, "Text 2", json.dumps(skills_completencies_dict), "", ""])
                        
                    except Exception as e:
                        print("Exception while offers data", e)
            # Save the file
            file_id = str(uuid4())
            wb.save(f"{file_id}.xlsx")
            prs = createPPT(file_id)
            prs.save(f"{file_id}.pptx")
            with open(f"{file_id}.pptx", "rb") as fh:
                response = HttpResponse(fh.read(), content_type="application/vnd.powerpoint")
                response["Content-Disposition"] = 'inline; filename="Insights_Book_Template.pptx"'
                return response
        # return JsonResponse({"mes":"success"})

        except Exception as e:
            return Response(
                data={"msg": "Unable to create and download ppt file."},
                content_type="application/json",
                status=status.HTTP_200_OK,
            )


# convert excel to ppt code


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_cell_border(cell, border_color="000000", border_width="12700"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    #     for lines in ['a:lnL','a:lnR','a:lnT','a:lnB']:
    lines = "a:lnB"
    ln = SubElement(tcPr, lines, w=border_width, cap="flat", cmpd="sng", algn="ctr")
    solidFill = SubElement(ln, "a:solidFill")
    srgbClr = SubElement(solidFill, "a:srgbClr", val=border_color)
    prstDash = SubElement(ln, "a:prstDash", val="solid")
    round_ = SubElement(ln, "a:round")
    headEnd = SubElement(ln, "a:headEnd", type="none", w="med", len="med")
    tailEnd = SubElement(ln, "a:tailEnd", type="none", w="med", len="med")


def createPPT(file_id):
    # Creating PPT file using excel data.
    prs = Presentation("sales2.0.pptx")

    # input_df=pd.read_excel('files/sales2.0_template.xlsx', keep_default_na=False,index_col=0)
    input_df = pd.read_excel(f"{file_id}.xlsx", keep_default_na=False, index_col=0)
    try:
        for slide_no, row in input_df.iterrows():
            slide = prs.slides[slide_no]
            for shape in slide.shapes:
                if shape.name.lower() == row.object_type.lower():
                    shape_name = shape.name
                    object_type = shape_name.split(" ")[0].lower()
                    if object_type == "title" or object_type == "subtitle" or object_type == "textbox":
                        tf = shape.text_frame
                        tf.clear()
                        tf.word_wrap = True
                        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        p = tf.paragraphs[0]
                        if row.font_size != "":
                            p.font.size = Pt(int(row.font_size))
                        if row.color != "":
                            run = p.add_run()
                            run.text = row.content
                            try:
                                color_code = row.color.split(",")
                                run.font.color.rgb = RGBColor(
                                    eval(color_code[0]), eval(color_code[1]), eval(color_code[2])
                                )
                            except:
                                print("Failed while applying color to font")
                        else:
                            p.text = row.content
                    elif object_type == "chart":
                        chart = shape.chart
                        chart_data = CategoryChartData()
                        input_data = json.loads(row.content)
                        chart_data.categories = input_data["categories"]
                        for key, value in input_data["series"].items():
                            chart_data.add_series(key, tuple(value))
                        chart.replace_data(chart_data)
                    elif object_type == "table":
                        tbl = shape._element.graphic.graphicData.tbl
                        if row.content != None:
                            for row_id, data in enumerate(eval(row.content)):
                                if row_id == 0:
                                    for col_id, row_data in enumerate(data):
                                        cell = shape.table.cell(row_id, col_id)
                                        _set_cell_border(cell)
                                        tf = cell.text_frame.paragraphs[0]
                                        tf.text = row_data
                                        if row.font_size != "":
                                            tf.font.size = Pt(row.font_size)
                                        tf.font.bold = True
                                else:
                                    tr = tbl.add_tr(height=50)
                                    for col_id, row_data in enumerate(data):
                                        tr.add_tc()
                                        cell = shape.table.cell(row_id, col_id)
                                        _set_cell_border(cell)
                                        tf = cell.text_frame.paragraphs[0]
                                        tf.text = str(row_data)
                                        tf.font.size = Pt(10)

                    elif object_type == "text":
                        try:
                            tf = shape.text_frame
                            tf.clear()
                            data_type = type(json.loads(row.content))
                            data = json.loads(row.content)
                            if data_type == list:
                                for list_value in range(0, len(data)):
                                    if list_value == 0:
                                        p = tf.paragraphs[0]
                                        p.text = data[list_value]
                                        p.level = 1
                                        p.font.size = Pt(10)
                                    else:
                                        p = tf.add_paragraph()
                                        p.text = data[list_value]
                                        p.level = 1
                                        p.font.size = Pt(10)

                            elif data_type == dict:
                                for key, value in data.items():
                                    p = tf.add_paragraph()
                                    p.text = key
                                    p.font.size = Pt(11)
                                    p.font.bold = True
                                    for value_list in range(len(value)):
                                        p = tf.add_paragraph()
                                        p.text = value[value_list]
                                        p.level = 1
                                        p.font.size = Pt(10)
                            else:
                                pass
                        except:
                            print("Exception while running text object")
                    elif object_type == "rectangle":
                        try:
                            left_pos = 0.6
                            top_pos = 2.1
                            data = json.loads(row.content)
                            offer_length = len(data)
                            for range_value in range(0, offer_length, 4):
                                diff_value = offer_length - range_value
                                if diff_value >= 4:
                                    end_value = range_value + 4
                                else:
                                    end_value = range_value + diff_value
                                for offer_no in range(range_value, end_value):
                                    txBox1 = slide.shapes.add_textbox(
                                        Inches(left_pos), Inches(top_pos), Inches(2.7), Inches(1.5)
                                    )
                                    tf = txBox1.text_frame
                                    tf.text = ""

                                    line = txBox1.line
                                    line.color.rgb = RGBColor(0, 0, 0)
                                    textbox_config = [
                                        {
                                            "left": 0.1,
                                            "top": 0.0,
                                            "width": 2.2,
                                            "height": 0.5,
                                            "bold": True,
                                            "font": 11,
                                            "text_color": False,
                                            "bg_color": False,
                                        },
                                        {
                                            "left": 0.1,
                                            "top": 1.0,
                                            "width": 1.0,
                                            "height": 0.5,
                                            "bold": False,
                                            "font": 14,
                                            "bg_color": False,
                                            "text_color": True,
                                        },
                                        {
                                            "left": 1.6,
                                            "top": 0.8,
                                            "width": 1.0,
                                            "height": 0.3,
                                            "bold": False,
                                            "font": 7,
                                            "bg_color": False,
                                            "text_color": False,
                                        },
                                        {
                                            "left": 1.7,
                                            "top": 1.0,
                                            "width": 0.8,
                                            "height": 0.3,
                                            "bold": False,
                                            "font": 10,
                                            "bg_color": True,
                                            "text_color": False,
                                        },
                                    ]
                                    shape_list = []
                                    for pos_no in range(len(textbox_config)):
                                        if pos_no == 0:
                                            text = data[offer_no]["offer"]
                                        if pos_no == 1:
                                            text = data[offer_no]["value"]
                                        if pos_no == 2:
                                            text = data[offer_no]["code"]
                                        if pos_no == 3:
                                            text = data[offer_no]["percent"]

                                        txBox = slide.shapes.add_textbox(
                                            Inches(left_pos + textbox_config[pos_no]["left"]),
                                            Inches(top_pos + textbox_config[pos_no]["top"]),
                                            Inches(textbox_config[pos_no]["width"]),
                                            Inches(textbox_config[pos_no]["height"]),
                                        )
                                        tf = txBox.text_frame
                                        tf.word_wrap = True
                                        p = tf.paragraphs[0]
                                        #                     p.text = add_text["text"]
                                        p.font.size = Pt(textbox_config[pos_no]["font"])
                                        if textbox_config[pos_no]["bold"]:
                                            p.font.bold = True
                                        if textbox_config[pos_no]["text_color"]:
                                            run = p.add_run()
                                            run.text = text
                                            run.font.color.rgb = RGBColor(0xFB, 0x8F, 0x00)
                                        else:
                                            p.text = text

                                        if textbox_config[pos_no]["bg_color"]:
                                            from pptx.enum.dml import MSO_THEME_COLOR

                                            # set fill type to solid color first
                                            txBox.fill.background()
                                            txBox.fill.solid()
                                            # set foreground (fill) color to a specific RGB color
                                            txBox.fill.fore_color.rgb = RGBColor(0xFB, 0x8F, 0x00)
                                            # change to a theme color
                                            txBox.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
                                            # set lighter or darker, -0.2 is 20% darker, 0.4 is 40% lighter
                                            txBox.fill.fore_color.brightness = 0.4
                                        shape_list.append(txBox)
                                    slide.shapes.add_group_shape(shapes=shape_list)
                                    left_pos = left_pos + 3.1
                                left_pos = 0.6
                                top_pos = top_pos + 1.6
                        except:
                            print("Something wrong in rectangle")
                    else:
                        print("There is no object type in Existed PPT file")
                        pass
        return prs
    except Exception as e:
        return Response(
            data={"msg": "Unable to download ppt file."},
            content_type="application/json",
            status=status.HTTP_200_OK,
        )
